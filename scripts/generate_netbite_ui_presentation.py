from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "deliverables" / "netbite-ui-presentation"
SCREEN_DIR = OUT_DIR / "screens"
OUT_FILE = OUT_DIR / "NetBite_UI_Presentation_Lazatin_Allen.pptx"

BG = "121014"
SURFACE = "1D191F"
SURFACE_2 = "262229"
BORDER = "3A3F3D"
TEXT = "DDD8DA"
MUTED = "9F9AA0"
RED = "D2474C"
ORANGE = "D18B5A"
SAGE = "79A99E"
BLUE = "6F96AD"
FONT = "Fira Code"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BG)
    for x in [0.42, 1.42, 2.42, 3.42, 4.42, 5.42, 6.42, 7.42, 8.42, 9.42, 10.42, 11.42, 12.42]:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), 0, Pt(0.4), Inches(7.5))
        line.fill.solid(); line.fill.fore_color.rgb = rgb("19171B"); line.line.fill.background()
    for y in [0.52, 1.52, 2.52, 3.52, 4.52, 5.52, 6.52]:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), Inches(13.333), Pt(0.4))
        line.fill.solid(); line.fill.fore_color.rgb = rgb("19171B"); line.line.fill.background()


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *, size: float = 18,
             color: str = TEXT, bold: bool = False, font: str = FONT,
             align: PP_ALIGN = PP_ALIGN.LEFT, valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
             spacing: float = 0.5, margin: float = 0.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.line_spacing = 1.05
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.kerning = Pt(spacing)
    return box


def add_rich_lines(slide, lines: list[tuple[str, str, bool]], x: float, y: float, w: float, h: float,
                   *, size: float = 16, gap: float = 7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    for index, (text, color, bold) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = text
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.1
        run = paragraph.runs[0]
        run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = rgb(color)
    return box


def add_panel(slide, x: float, y: float, w: float, h: float, *, border: str = BORDER,
              fill: str = SURFACE, rail: str | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(border); shape.line.width = Pt(0.8)
    if rail:
        rail_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.055), Inches(h))
        rail_shape.fill.solid(); rail_shape.fill.fore_color.rgb = rgb(rail); rail_shape.line.fill.background()
    return shape


def add_header(slide, number: str, kicker: str, title: str, subtitle: str | None = None) -> None:
    add_text(slide, number, 0.58, 0.38, 0.55, 0.25, size=10, color=RED, bold=True)
    add_text(slide, kicker.upper(), 1.16, 0.37, 4.4, 0.28, size=10, color=ORANGE, bold=True)
    add_text(slide, title.upper(), 0.58, 0.78, 12.1, 0.55, size=26, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.31, 11.8, 0.35, size=12.5, color=MUTED)
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.72), Inches(12.15), Pt(1))
    divider.fill.solid(); divider.fill.fore_color.rgb = rgb(BORDER); divider.line.fill.background()


def add_footer(slide, number: int) -> None:
    add_text(slide, "NETBITE / UI OVERVIEW", 0.58, 7.12, 3.2, 0.18, size=8.5, color=MUTED)
    add_text(slide, f"{number:02d}", 12.08, 7.1, 0.62, 0.2, size=9, color=RED, bold=True, align=PP_ALIGN.RIGHT)


def add_screenshot(slide, name: str, x: float, y: float, w: float, h: float, alt: str) -> None:
    add_panel(slide, x - 0.08, y - 0.08, w + 0.16, h + 0.16, border=BORDER, fill=SURFACE_2)
    picture = slide.shapes.add_picture(str(SCREEN_DIR / name), Inches(x), Inches(y), Inches(w), Inches(h))
    picture._element.nvPicPr.cNvPr.set("descr", alt)


def add_tag(slide, text: str, x: float, y: float, w: float, *, tone: str = RED) -> None:
    add_panel(slide, x, y, w, 0.34, border=tone, fill=BG)
    add_text(slide, text.upper(), x + 0.1, y + 0.085, w - 0.2, 0.16, size=8.5, color=tone, bold=True, align=PP_ALIGN.CENTER)


def build_deck() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 01 — Title
    slide = prs.slides.add_slide(blank); set_background(slide)
    logo = ROOT / "assets" / "images" / "branding" / "netbite-menu-logo-mobile.png"
    slide.shapes.add_picture(str(logo), Inches(0.78), Inches(0.66), Inches(1.15), Inches(1.15))
    add_text(slide, "NETBITE", 0.78, 2.05, 6.7, 0.78, size=39, color=TEXT, bold=True)
    add_text(slide, "MOBILE NETWORKING EDUCATION", 0.82, 2.82, 5.9, 0.34, size=13, color=ORANGE, bold=True)
    add_text(slide, "Learn networking concepts through concise lessons, accurate visuals, guided labs, and a deterministic network sandbox.", 0.82, 3.42, 5.65, 1.12, size=18, color=TEXT)
    add_tag(slide, "LEARN", 0.82, 4.92, 1.35, tone=RED)
    add_tag(slide, "PRACTICE", 2.32, 4.92, 1.62, tone=ORANGE)
    add_tag(slide, "EXPERIMENT", 4.08, 4.92, 1.85, tone=SAGE)
    add_text(slide, "UI PRESENTATION", 0.82, 6.52, 2.2, 0.25, size=10, color=MUTED, bold=True)
    add_text(slide, "ALLEN LAZATIN / ITE-231", 0.82, 6.83, 3.4, 0.25, size=10, color=TEXT)
    # Device topology
    pc = ROOT / "assets" / "images" / "devices" / "device-pc-mobile.png"
    switch = ROOT / "assets" / "images" / "devices" / "device-switch-mobile.png"
    router = ROOT / "assets" / "images" / "devices" / "device-router-mobile.png"
    add_panel(slide, 7.25, 0.72, 5.45, 5.98, border=BORDER, fill=SURFACE, rail=RED)
    add_text(slide, "GUIDED NETWORK PATH", 7.72, 1.15, 3.5, 0.25, size=10, color=SAGE, bold=True)
    for x1, y1, x2, y2, color in [(8.55, 3.15, 9.55, 3.15, RED), (10.17, 3.15, 11.18, 3.15, ORANGE)]:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y1), Inches(x2-x1), Pt(2.2))
        line.fill.solid(); line.fill.fore_color.rgb = rgb(color); line.line.fill.background()
    slide.shapes.add_picture(str(pc), Inches(7.65), Inches(2.48), Inches(1.25), Inches(1.25))
    slide.shapes.add_picture(str(switch), Inches(9.33), Inches(2.5), Inches(1.32), Inches(1.32))
    slide.shapes.add_picture(str(router), Inches(11.08), Inches(2.48), Inches(1.28), Inches(1.28))
    add_text(slide, "PC", 7.95, 3.88, 0.68, 0.2, size=9, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "SWITCH", 9.56, 3.88, 0.88, 0.2, size=9, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "ROUTER", 11.27, 3.88, 0.92, 0.2, size=9, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "A restrained retro-console interface keeps the focus on decisions, relationships, and feedback.", 7.72, 4.65, 4.45, 0.86, size=15, color=MUTED)
    add_text(slide, "STATE-BASED EDUCATIONAL SIMULATION / NO LIVE PACKETS", 7.72, 5.91, 4.38, 0.28, size=8.5, color=ORANGE, bold=True)

    # 02 — About
    slide = prs.slides.add_slide(blank); set_background(slide)
    add_header(slide, "02", "Product overview", "What is NetBite?", "A beginner-focused mobile learning system for foundational computer networking.")
    add_panel(slide, 0.58, 2.08, 5.55, 3.95, border=BORDER, fill=SURFACE, rail=RED)
    add_text(slide, "THE LEARNING PROBLEM", 0.9, 2.42, 4.8, 0.25, size=11, color=RED, bold=True)
    add_text(slide, "Definitions alone can hide the decisions networking devices make.", 0.9, 2.86, 4.75, 0.92, size=17.5, bold=True)
    add_rich_lines(slide, [
        ("• Abstract processes are hard to visualize.", TEXT, False),
        ("• Large simulators can overwhelm beginners.", TEXT, False),
        ("• Wrong answers need explanations, not only scores.", TEXT, False),
    ], 0.92, 4.0, 4.75, 1.55, size=13.8, gap=8)
    add_panel(slide, 6.38, 2.08, 6.35, 3.95, border=SAGE, fill=SURFACE, rail=SAGE)
    add_text(slide, "THE NETBITE RESPONSE", 6.72, 2.42, 5.4, 0.25, size=11, color=SAGE, bold=True)
    add_rich_lines(slide, [
        ("12 CHAPTERS", ORANGE, True),
        ("84 short, focused lessons", TEXT, False),
        ("Guided labs, CLI practice, quizzes, and active-recall flashcards", TEXT, False),
        ("A bounded sandbox that explains what the modeled result proves", TEXT, False),
    ], 6.72, 2.9, 5.35, 2.35, size=16, gap=8)
    add_text(slide, "OFFLINE-FIRST / MOBILE-FIRST / DETERMINISTIC", 6.72, 5.5, 5.05, 0.26, size=9.5, color=ORANGE, bold=True)
    add_footer(slide, 2)

    # 03 — Welcome
    slide = prs.slides.add_slide(blank); set_background(slide)
    add_header(slide, "03", "First launch", "A clear way into the app", "Accounts are useful, but never block a learner from starting locally.")
    add_screenshot(slide, "01-account-welcome.png", 0.72, 1.96, 3.35, 5.23, "NetBite first-launch account welcome screen")
    add_text(slide, "THREE RECOGNIZABLE CHOICES", 4.55, 2.18, 5.7, 0.28, size=11, color=ORANGE, bold=True)
    add_text(slide, "The first screen explains the tradeoff before asking the learner to decide.", 4.55, 2.62, 7.25, 0.75, size=22, bold=True)
    add_rich_lines(slide, [
        ("SIGN IN — restore cloud progress", RED, True),
        ("CREATE ACCOUNT — begin optional backup", TEXT, False),
        ("CONTINUE AS GUEST — learn immediately on the device", SAGE, True),
    ], 4.55, 3.66, 7.25, 1.6, size=16, gap=10)
    add_panel(slide, 4.55, 5.68, 7.25, 0.93, border=BORDER, fill=SURFACE_2, rail=ORANGE)
    add_text(slide, "HCI PRINCIPLE", 4.84, 5.9, 1.55, 0.22, size=9, color=ORANGE, bold=True)
    add_text(slide, "Clear alternatives reduce uncertainty and preserve learner control.", 6.32, 5.84, 4.9, 0.38, size=13.5, color=TEXT)
    add_footer(slide, 3)

    # 04 — Main menu and learning path
    slide = prs.slides.add_slide(blank); set_background(slide)
    add_header(slide, "04", "Navigation hierarchy", "Learning stays the primary action", "The interface promotes the next useful step without hiding exploration tools.")
    add_screenshot(slide, "02-main-menu.png", 0.7, 1.96, 3.2, 5.0, "NetBite main menu with Start Learning as the dominant action")
    add_screenshot(slide, "03-learning-path.png", 4.2, 1.96, 3.2, 5.0, "NetBite learning path with chapter progress and saved learning")
    add_rich_lines(slide, [
        ("ONE DOMINANT ACTION", RED, True),
        ("Continue the next unfinished lesson without searching.", TEXT, False),
        ("RECOGNITION OVER RECALL", ORANGE, True),
        ("Chapters, progress, review, and saved items use consistent labels and icons.", TEXT, False),
        ("TOOLS STAY VISIBLE", SAGE, True),
        ("The Sandbox is distinct from the guided curriculum but remains easy to find.", TEXT, False),
    ], 7.86, 2.15, 4.55, 3.75, size=14.2, gap=7)
    add_footer(slide, 4)

    # 05 — Lesson UI
    slide = prs.slides.add_slide(blank); set_background(slide)
    add_header(slide, "05", "Lesson experience", "Making invisible protocols visible", "Exact values remain selectable and responsive instead of being baked into an image.")
    add_screenshot(slide, "04-arp-lesson.png", 0.7, 1.96, 3.72, 5.0, "ARP lesson showing outer Ethernet frame and ARP request fields")
    add_text(slide, "ARP REQUEST / OPERATIONAL VIEW", 4.9, 2.12, 6.8, 0.3, size=11, color=ORANGE, bold=True)
    add_text(slide, "The lesson reveals what is sent, what each device inspects, and where the behavior stops.", 4.9, 2.55, 7.15, 0.8, size=21, bold=True)
    add_rich_lines(slide, [
        ("FF:FF:FF:FF:FF:FF", RED, True),
        ("Ethernet broadcast destination", MUTED, False),
        ("0x0806", ORANGE, True),
        ("EtherType identifying ARP", MUTED, False),
        ("UNKNOWN / UNUSED", SAGE, True),
        ("ARP target-hardware field before the mapping is known", MUTED, False),
    ], 4.92, 3.72, 6.2, 2.2, size=14.5, gap=4)
    add_text(slide, "FIELDS STACK ON MOBILE INSTEAD OF SHRINKING", 4.92, 6.36, 6.4, 0.24, size=9.5, color=SAGE, bold=True)
    add_footer(slide, 5)

    # 06 — Sandbox / close
    slide = prs.slides.add_slide(blank); set_background(slide)
    add_header(slide, "06", "Practice and experimentation", "From guided learning to safe tinkering", "The Network Sandbox turns configuration into an explainable state-based result.")
    add_screenshot(slide, "05-network-sandbox.png", 0.7, 1.96, 3.72, 5.0, "NetBite routed network sandbox guide")
    add_text(slide, "A COMPLETE ROUTED PRESET", 4.9, 2.1, 5.7, 0.3, size=11, color=ORANGE, bold=True)
    add_text(slide, "Learners can inspect, test, break, and restore a working network without risking real equipment.", 4.9, 2.52, 7.15, 0.82, size=21, bold=True)
    steps = [
        ("01", "INSPECT", "See addressing, interfaces, gateways, VLANs, and routes."),
        ("02", "TEST", "Send a frame or run ping through deterministic forwarding logic."),
        ("03", "EXPLAIN", "Read what happened, what it proves, and the next useful check."),
        ("04", "RECOVER", "Use hints, undo, reset, or a known-good preset."),
    ]
    y = 3.72
    for number, label, detail in steps:
        add_text(slide, number, 4.92, y, 0.42, 0.22, size=10, color=RED, bold=True)
        add_text(slide, label, 5.45, y - 0.02, 1.35, 0.24, size=11, color=ORANGE, bold=True)
        add_text(slide, detail, 6.78, y - 0.04, 5.1, 0.43, size=12.5, color=TEXT)
        y += 0.68
    add_text(slide, "NETBITE / LEARN THE RULE — THEN TEST THE RULE", 4.92, 6.55, 6.75, 0.25, size=10, color=SAGE, bold=True)
    add_footer(slide, 6)

    prs.core_properties.title = "NetBite UI Presentation"
    prs.core_properties.subject = "Mobile networking education application UI overview"
    prs.core_properties.author = "Allen Lazatin"
    prs.core_properties.keywords = "NetBite, networking, mobile, education, UI, simulation"
    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    build_deck()
